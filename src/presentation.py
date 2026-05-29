import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pathlib import Path
import logging

logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s [%(levelname)s] %(message)s"
)
logger = logging.getLogger(__name__)

# Cores padrão
AZUL_ESCURO  = RGBColor(0x1F, 0x4E, 0x79)
AZUL_MEDIO   = RGBColor(0x2E, 0x75, 0xB6)
AZUL_CLARO   = RGBColor(0xA8, 0xC8, 0xE8)
BRANCO       = RGBColor(0xFF, 0xFF, 0xFF)
CINZA        = RGBColor(0x59, 0x59, 0x59)
VERMELHO     = RGBColor(0xC0, 0x00, 0x00)

IMG_DIR    = Path("output/presentation")
OUTPUT_PPT = Path("output/presentation/RAIS_Insights_2022.pptx")


# ── Helpers ──────────────────────────────────────────────────────────────────

def add_text(tf, text, size, bold=False, color=BRANCO, align=PP_ALIGN.LEFT):
    tf.text = ""
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def set_bg(slide, color: RGBColor):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_box(slide, text, left, top, width, height,
            font_size=12, bold=False, font_color=BRANCO,
            bg_color=AZUL_ESCURO, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    box.fill.solid()
    box.fill.fore_color.rgb = bg_color
    tf = box.text_frame
    tf.word_wrap = True
    add_text(tf, text, font_size, bold, font_color, align)
    return box


# ── Slides ────────────────────────────────────────────────────────────────────

def slide_capa(prs, df):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, AZUL_ESCURO)

    # Faixa decorativa
    bar = slide.shapes.add_shape(1,
        Inches(0), Inches(2.8), Inches(10), Inches(0.08))
    bar.fill.solid()
    bar.fill.fore_color.rgb = AZUL_CLARO
    bar.line.fill.background()

    add_box(slide, "Análise do Mercado de Trabalho Formal",
            0.4, 1.0, 9.2, 1.0, font_size=28, bold=True, align=PP_ALIGN.CENTER)

    add_box(slide, "Microdados RAIS 2022 — Simulado com base na estrutura oficial do MTE",
            0.4, 2.1, 9.2, 0.6, font_size=14, align=PP_ALIGN.CENTER)

    total   = df.shape[0]
    salario = df["SALARIO_MENSAL"].mean()
    ativos  = (df["SITUACAO_VINCULO"] == "Ativo").sum()

    metricas = [
        (f"{total:,}", "Vínculos\nAnalisados"),
        (f"R$ {salario:,.0f}", "Salário\nMédio"),
        (f"{ativos:,}", "Vínculos\nAtivos"),
    ]
    for i, (val, label) in enumerate(metricas):
        x = 0.5 + i * 3.1
        add_box(slide, val,  x, 3.2, 2.8, 0.6,
                font_size=22, bold=True, align=PP_ALIGN.CENTER)
        add_box(slide, label, x, 3.85, 2.8, 0.55,
                font_size=11, bg_color=AZUL_MEDIO, align=PP_ALIGN.CENTER)

    add_box(slide, "Projeto 1 — Portfólio de Engenharia de Dados",
            0.4, 6.6, 9.2, 0.4, font_size=10,
            bg_color=RGBColor(0x17, 0x37, 0x59), align=PP_ALIGN.CENTER)

    logger.info("Slide 1: Capa criada")


def slide_metodologia(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BRANCO)

    add_box(slide, "Metodologia", 0.4, 0.3, 9.2, 0.6,
            font_size=22, bold=True, font_color=AZUL_ESCURO,
            bg_color=BRANCO, align=PP_ALIGN.LEFT)

    bar = slide.shapes.add_shape(1,
        Inches(0.4), Inches(1.0), Inches(9.2), Inches(0.05))
    bar.fill.solid()
    bar.fill.fore_color.rgb = AZUL_MEDIO
    bar.line.fill.background()

    etapas = [
        ("📥", "Ingestão",       "Leitura do Excel com\nPandas + inspeção\ninicial dos dados"),
        ("🧹", "Limpeza",        "Tratamento de nulos,\nvalores inválidos e\npadronização de tipos"),
        ("📊", "Análise",        "Geração de métricas,\ninsights e 7 gráficos\ncom Matplotlib/Seaborn"),
        ("📑", "Apresentação",   "Slides automáticos\ngerados com\npython-pptx"),
    ]

    for i, (icon, titulo, desc) in enumerate(etapas):
        x = 0.4 + i * 2.35
        add_box(slide, icon,   x, 1.3,  2.1, 0.55,
                font_size=24, align=PP_ALIGN.CENTER)
        add_box(slide, titulo, x, 1.9,  2.1, 0.45,
                font_size=13, bold=True, bg_color=AZUL_MEDIO, align=PP_ALIGN.CENTER)
        add_box(slide, desc,   x, 2.45, 2.1, 1.0,
                font_size=10, font_color=CINZA, bg_color=BRANCO, align=PP_ALIGN.CENTER)

        if i < 3:
            arr = slide.shapes.add_textbox(
                Inches(x + 2.1), Inches(1.9), Inches(0.25), Inches(0.45))
            arr.text_frame.text = "→"
            arr.text_frame.paragraphs[0].runs[0].font.size = Pt(18)
            arr.text_frame.paragraphs[0].runs[0].font.color.rgb = AZUL_MEDIO

    stack = [
        "🐍 Python 3.11",
        "🐼 Pandas",
        "📈 Matplotlib / Seaborn",
        "📊 OpenPyXL",
        "📑 python-pptx",
    ]
    add_box(slide, "Stack Utilizada",
            0.4, 3.7, 9.2, 0.4,
            font_size=13, bold=True, bg_color=AZUL_ESCURO, align=PP_ALIGN.CENTER)
    add_box(slide, "   ".join(stack),
            0.4, 4.2, 9.2, 0.5,
            font_size=11, font_color=CINZA, bg_color=BRANCO, align=PP_ALIGN.CENTER)

    logger.info("Slide 2: Metodologia criada")


def slide_grafico(prs, titulo, subtitulo, img_file, insight):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, BRANCO)

    add_box(slide, titulo, 0.4, 0.15, 9.2, 0.5,
            font_size=18, bold=True, font_color=AZUL_ESCURO,
            bg_color=BRANCO, align=PP_ALIGN.LEFT)
    add_box(slide, subtitulo, 0.4, 0.65, 9.2, 0.35,
            font_size=11, font_color=CINZA,
            bg_color=BRANCO, align=PP_ALIGN.LEFT)

    img_path = IMG_DIR / img_file
    if img_path.exists():
        slide.shapes.add_picture(
            str(img_path),
            Inches(0.4), Inches(1.1),
            Inches(9.2), Inches(4.6)
        )

    add_box(slide, f"💡  {insight}",
            0.4, 5.85, 9.2, 0.65,
            font_size=11, bg_color=AZUL_ESCURO, align=PP_ALIGN.LEFT)

    return slide


def slide_conclusao(prs, df):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide, AZUL_ESCURO)

    add_box(slide, "Conclusões & Próximos Passos",
            0.4, 0.3, 9.2, 0.6,
            font_size=22, bold=True, align=PP_ALIGN.CENTER)

    bar = slide.shapes.add_shape(1,
        Inches(0.4), Inches(1.0), Inches(9.2), Inches(0.05))
    bar.fill.solid()
    bar.fill.fore_color.rgb = AZUL_CLARO
    bar.line.fill.background()

    sal_masc = df[df["GENERO"] == "Masculino"]["SALARIO_MENSAL"].mean()
    sal_fem  = df[df["GENERO"] == "Feminino"]["SALARIO_MENSAL"].mean()
    gap      = abs(((sal_masc - sal_fem) / sal_masc) * 100)
    top_setor = df.groupby("SETOR_ATIVIDADE")["SALARIO_MENSAL"].mean().idxmax()
    retencao  = (df["SITUACAO_VINCULO"] == "Ativo").sum() / len(df) * 100

    conclusoes = [
        f"🏆  {top_setor} lidera os maiores salários médios da base",
        f"⚖️   Gap salarial de gênero de {gap:.1f}% — tema relevante para RH",
        f"📋  Taxa de retenção de {retencao:.1f}% indica baixa rotatividade",
        f"🎓  Escolaridade superior está associada a maiores faixas salariais",
        f"👥  Faixa etária 35-54 anos concentra a maior parte dos vínculos",
    ]

    for i, texto in enumerate(conclusoes):
        add_box(slide, texto,
                0.5, 1.3 + i * 0.85, 9.0, 0.65,
                font_size=12, bg_color=AZUL_MEDIO, align=PP_ALIGN.LEFT)

    add_box(slide, "Próximo Projeto →  Banco de dados PostgreSQL + Power BI",
            0.4, 5.8, 9.2, 0.55,
            font_size=12, bold=True,
            bg_color=RGBColor(0x17, 0x37, 0x59), align=PP_ALIGN.CENTER)

    logger.info("Slide final: Conclusões criadas")


# ── Main ──────────────────────────────────────────────────────────────────────

def gerar_apresentacao(caminho_dados: str) -> None:
    logger.info("Iniciando geração da apresentação...")
    df = pd.read_excel(caminho_dados)

    prs = Presentation()
    prs.slide_width  = Inches(10)
    prs.slide_height = Inches(7.5)

    slide_capa(prs, df)
    slide_metodologia(prs)

    graficos = [
        ("Salário Médio por Setor",
         "Comparativo entre os 10 setores de atividade econômica",
         "01_salario_por_setor.png",
         "Tecnologia da Informação lidera com maior salário médio, seguida de Serviços Financeiros e Saúde"),

        ("Distribuição por Gênero",
         "Proporção de vínculos masculinos e femininos na base",
         "02_distribuicao_genero.png",
         "Distribuição equilibrada entre gêneros, com leve predominância masculina"),

        ("Distribuição por Escolaridade",
         "Nível de escolaridade dos trabalhadores formais",
         "03_escolaridade.png",
         "Ensino Médio Completo é o nível mais frequente, seguido de Superior Completo"),

        ("Distribuição por Faixa Etária",
         "Concentração etária dos vínculos empregatícios",
         "04_faixa_etaria.png",
         "Trabalhadores entre 35 e 54 anos representam a maior parcela da força de trabalho formal"),

        ("Gap Salarial por Gênero e Setor",
         "Comparativo do salário médio entre homens e mulheres por setor",
         "05_salario_genero_setor.png",
         "Disparidades salariais de gênero variam por setor — TI e Finanças apresentam maiores gaps"),

        ("Situação dos Vínculos",
         "Proporção entre vínculos ativos e desligados em 2022",
         "06_situacao_vinculo.png",
         "83,5% dos vínculos estavam ativos ao final de 2022, indicando baixa rotatividade"),

        ("Top 10 Cargos Mais Frequentes",
         "Cargos com maior volume de vínculos formais na base",
         "07_top_cargos.png",
         "Cargos operacionais e de saúde dominam o ranking de frequência"),
    ]

    for titulo, subtitulo, img, insight in graficos:
        slide_grafico(prs, titulo, subtitulo, img, insight)

    slide_conclusao(prs, df)

    prs.save(OUTPUT_PPT)
    logger.info(f"✅ Apresentação salva em: {OUTPUT_PPT}")
    logger.info(f"   Total de slides: {len(prs.slides)}")


if __name__ == "__main__":
    gerar_apresentacao("data/processed/RAIS_tratado.xlsx")